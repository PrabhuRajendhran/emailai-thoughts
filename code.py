from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE

# Create a presentation object
prs = Presentation()

# Slide 1: Title Slide
slide = prs.slides.add_slide(prs.slide_layouts[0])
slide.shapes.title.text = "Enterprise Content Routing: From Team-Scale LLMs to Future-Proof Hybrid AI"
slide.placeholders[1].text = "Scalable, Auditable, Multi-Content, Enterprise-Ready"

# Slide 2: LLM vs Hybrid Table
slide = prs.slides.add_slide(prs.slide_layouts[5])
shapes = slide.shapes
shapes.title.text = "LLM vs Hybrid – Side-by-Side Comparison"

rows = 10
cols = 3
left = Inches(0.5)
top = Inches(1.5)
width = Inches(9)
height = Inches(3)
table = shapes.add_table(rows, cols, left, top, width, height).table

# Set column headers
table.cell(0,0).text = "Dimension"
table.cell(0,1).text = "LLM-Prompt"
table.cell(0,2).text = "Hybrid Platform"

# Fill table rows
data = [
    ("Scalability Across Units","Single model bottleneck","Modular + thresholds per unit"),
    ("Cost Efficiency","Very high at scale","Low–moderate: embeddings + re-rankers"),
    ("Latency / SLA","Risk >5s","Deterministic; meets 3–5s SLA"),
    ("Governance & Audit","Weak; generative output","Strong; top-K, confidence, human-in-loop"),
    ("Taxonomy Handling","Messy categories confuse output","Hybrid keyword+semantic ranking; thresholds"),
    ("Unit-Level Adaptation","Expensive, complex","Per-unit thresholds; optional fine-tuning"),
    ("Compliance Fit","Weak; hard to explain","Deterministic, auditable, escalation rules"),
    ("Document Type Flexibility","Primarily email","Email, PDFs, OCR, chat, tickets"),
    ("Long-Term Value","Point solution","Enterprise-grade, reusable, future-proof")
]

for i, (d, l, h) in enumerate(data, start=1):
    table.cell(i,0).text = d
    table.cell(i,1).text = l
    table.cell(i,2).text = h

# Slide 3: Key Strategic Advantages
slide = prs.slides.add_slide(prs.slide_layouts[5])
shapes = slide.shapes
shapes.title.text = "Key Strategic Advantages"
# For simplicity, add text boxes for icons
advantages = ["Predictable Cost & Latency","Governance & Audit","Unit-Level Customization","Multi-Content Flexibility","Future-Proof Infrastructure"]
for i, adv in enumerate(advantages):
    left = Inches(0.5 + i*1.8)
    top = Inches(1.5)
    width = Inches(1.5)
    height = Inches(1.0)
    box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    box.text = adv

# Slide 4: Pipeline Diagram (simplified)
slide = prs.slides.add_slide(prs.slide_layouts[5])
shapes = slide.shapes
shapes.title.text = "Hybrid Platform Pipeline (Top-K=100→5→2→1)"
steps = ["Input: Email/PDF/Chat/Ticket", "Embedding + Top-K Retrieval (100→5)", "Cross-Encoder Re-Ranker (5→2)", "Threshold Re-Ranker (2→1)", "Routing/Escalation"]
for i, step in enumerate(steps):
    left = Inches(0.5)
    top = Inches(1.0 + i*0.7)
    width = Inches(9)
    height = Inches(0.6)
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    box.text = step

# Slide 5: ROI & Strategic Value
slide = prs.slides.add_slide(prs.slide_layouts[5])
shapes.title.text = "ROI & Strategic Value"
boxes = ["Reduced Operational Cost","Compliance Risk Mitigation","Enterprise Scalability","Future-Proof Infrastructure"]
for i, b in enumerate(boxes):
    left = Inches(0.5 + i*2.0)
    top = Inches(1.5)
    width = Inches(1.8)
    height = Inches(1.0)
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.text = b

# Slide 6: Cheat Sheet (bullet points)
slide = prs.slides.add_slide(prs.slide_layouts[1])
slide.shapes.title.text = "MD/Board Cheat Sheet"
body = slide.placeholders[1]
body.text = "- Scalability: Modular core, per-unit thresholds\n- Cost: Embeddings + re-rankers → low per-email cost\n- Latency/SLA: ~4 sec/email\n- Governance & Audit: Top-K, confidence scores, human-in-loop\n- Taxonomy Handling: Hybrid keyword + semantic ranking; thresholds\n- Unit Adaptation: Thresholds per unit; optional fine-tuning\n- Compliance Fit: Deterministic, auditable, escalation rules\n- Document Types: Email, PDF, OCR, chat, tickets\n- Long-Term Value: Enterprise-grade, reusable, future-proof"

# Slide 7: Architecture Options Comparison (simplified)
slide = prs.slides.add_slide(prs.slide_layouts[5])
shapes.title.text = "Architecture Options Comparison"
body_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(9), Inches(3))
body_box.text = "Color-coded table highlighting LLM, SLM, Embedder+NN, Embedder+Re-ranker, Hybrid Embedder+Re-ranker. Hybrid single-pipeline is optimal."

# Slide 8: Infrastructure & Deployment
slide = prs.slides.add_slide(prs.slide_layouts[5])
shapes.title.text = "Infrastructure & Deployment"
body_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(9), Inches(3))
body_box.text = "- GPU Cluster (2–4 GPUs) → Cross-Encoder\n- CPU Nodes → Embedding + Pre/Post-processing\n- Load Balancer → Multi-unit pipelines\n- Optional Unit-Specific Threshold Layer"

# Save file
file_path = "/mnt/data/Enterprise_Hybrid_AI_Deck_Final.pptx"
prs.save(file_path)
